"""Generate IBR Technical Deep-Dive: Template Creation & Hybrid Network.

Author: Supreeth Mysore, Lead GPU Architect
"""
import os, sys
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from create_presentation import (
    PRIMARY, PRIMARY_LIGHT, PRIMARY_DARK, WHITE, TINT, DARK, BODY, MUTED, BORDER,
    SLIDE_W, SLIDE_H,
    add_badge, add_page_num, add_title, add_subtitle, add_insight_bar,
    add_body_text, add_stat_box, add_card, add_process_step, dark_slide, light_slide,
)

import create_presentation as cp
cp.TOTAL_SLIDES = 30
TOTAL_SLIDES = 30


def _dark_title(s, num, title, sub):
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = num; r.font.size = Pt(60); r.font.bold = True; r.font.color.rgb = PRIMARY_LIGHT; r.font.name = "Georgia"
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = title; r2.font.size = Pt(32); r2.font.bold = True; r2.font.color.rgb = WHITE; r2.font.name = "Georgia"
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = PRIMARY_LIGHT; ln.line.fill.background()
    t3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    r3 = t3.text_frame.paragraphs[0].add_run()
    r3.text = sub; r3.font.size = Pt(14); r3.font.color.rgb = RGBColor(0xCC,0xCC,0xCC); r3.font.name = "Arial"


def _closing_line(s, text, y):
    t = s.shapes.add_textbox(Inches(0.8), Inches(y), Inches(8.4), Inches(0.3))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = text; r.font.size = Pt(12); r.font.color.rgb = RGBColor(0xCC,0xCC,0xCC); r.font.name = "Arial"


def _add_table(slide, headers, rows, x=0.5, y=1.4, w=9, h=2.2):
    t = slide.shapes.add_table(len(rows)+1, len(headers), Inches(x), Inches(y), Inches(w), Inches(h))
    tbl = t.table
    for ci, hdr in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = hdr; c.fill.solid(); c.fill.fore_color.rgb = PRIMARY
        for p in c.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                r.font.color.rgb = WHITE; r.font.bold = True; r.font.size = Pt(10); r.font.name = "Arial"
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            c = tbl.cell(ri+1, ci)
            c.text = val
            if ri % 2 == 0: c.fill.solid(); c.fill.fore_color.rgb = TINT
            for p in c.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for r in p.runs:
                    r.font.size = Pt(10); r.font.name = "Arial"; r.font.color.rgb = DARK


def build():
    pres = Presentation()
    pres.slide_width = SLIDE_W; pres.slide_height = SLIDE_H
    n = [0]
    def sn(): n[0] += 1; return n[0]

    # S1: TITLE
    s = dark_slide(pres); sn()
    t = s.shapes.add_textbox(Inches(0.8), Inches(0.9), Inches(8.4), Inches(0.8))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "IBR Template Creation &\nHybrid Network Architecture"; r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Georgia"
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.4), Inches(0.4))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = "Technical Deep-Dive: Grid-Mesh Pipeline, Results Walkthrough, FastAPI Dashboard"; r2.font.size = Pt(14); r2.font.color.rgb = PRIMARY_LIGHT; r2.font.name = "Arial"
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.45), Inches(3), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = PRIMARY_LIGHT; ln.line.fill.background()
    for i, txt in enumerate(["Supreeth Mysore  |  Lead GPU Architect", "Hitachi Digital Services  |  Pratt & Whitney F135 Program", "March 2026  |  Sprint 3 Technical Review"]):
        _closing_line(s, txt, 2.7 + i*0.35)

    # S2: AGENDA
    s = light_slide(pres); pn = sn()
    add_badge(s, "OVERVIEW"); add_page_num(s, pn); add_title(s, "Technical Agenda"); add_subtitle(s, "Deep-dive into how the template is built and how every result traces back to it")
    for i, (t, d) in enumerate([("1. Template Architecture: Grid + Mesh Hybrid", "Why two data structures; how they complement each other"),
        ("2. Voxel Grid: Discretizing 3D Space", "The code that turns 55K raw points into 22K indexed cells"),
        ("3. KD-Tree Mesh Network", "Spatial queries at O(log n) — the backbone of every phase"),
        ("4. Zone Template: 13-Zone Mapping", "Percentage boundaries, multi-zone resolution, compliance logic"),
        ("5. Results Trace: End-to-End Defect Walk", "Following one defect through all 8 phases to its disposition"),
        ("6. FastAPI Dashboard: Code Walkthrough", "12 API endpoints, Plotly 3D rendering, live pipeline execution")]):
        add_card(s, t, d, 0.5, 1.35 + i*0.47, w=9, h=0.42)
    add_insight_bar(s, "THIS DECK", "Goes deep into HOW the template is built, HOW results are computed, and HOW the dashboard works. Code snippets included.")

    # S3: SECTION 01
    s = dark_slide(pres); sn()
    _dark_title(s, "01", "Template Architecture", "How we build the computational substrate for defect detection")

    # S4: WHY HYBRID
    s = light_slide(pres); pn = sn()
    add_badge(s, "ARCHITECTURE"); add_page_num(s, pn); add_title(s, "Why Hybrid Grid + Mesh?"); add_subtitle(s, "Two data structures, each optimized for different query patterns")
    add_card(s, "Voxel Grid (Discretization)", "O(1) lookup via floor(point / voxel_size)\nUniform spatial partitioning\n55K raw points \u2192 22K unique cells at 0.5mm\nUsed for: downsampling, zone lookup", 0.5, 1.4, w=4.2, h=1.5)
    add_card(s, "KD-Tree Mesh (Spatial Queries)", "O(log n) nearest-neighbor queries\nscipy.spatial.cKDTree (C-compiled)\n22K queries in <1 second with workers=-1\nUsed for: deviation, outlier, adjacency", 5.2, 1.4, w=4.2, h=1.5)
    add_body_text(s, "Grid alone cannot do nearest-neighbor queries on irregular geometry.\nKD-tree alone cannot do O(1) cell-based zone lookup.\nTogether they cover every query pattern in the pipeline.", x=0.5, y=3.15, h=0.8, size=11, color=BODY)
    add_insight_bar(s, "DESIGN RATIONALE", "Grid gives speed for filtering. KD-tree gives precision for measurement. Neither alone is sufficient.")

    # S5: 5-STEP TEMPLATE CREATION
    s = light_slide(pres); pn = sn()
    add_badge(s, "TEMPLATE"); add_page_num(s, pn); add_title(s, "Template Creation: 5-Step Pipeline"); add_subtitle(s, "From raw PLY to indexed, normal-enriched point cloud")
    for i, (num, title, desc) in enumerate([("1","PLY Load","55K raw\npoints"), ("2","Voxel Grid","floor(pt/0.5mm)\n22K cells"), ("3","Outlier Filter","KD-tree k=20\n\u03bc+2.5\u03c3 threshold"), ("4","Normals","PCA k=30\nsmallest eigvec"), ("5","Template","22,064 pts\nindexed+normals")]):
        add_process_step(s, num, title, desc, 0.5 + i*1.8, 1.5)
    add_body_text(s, "The template IS the cleaned point cloud with normals.\nEvery downstream phase queries this structure via grid index or KD-tree.", x=0.5, y=3.2, h=0.6, size=11, color=BODY)
    add_insight_bar(s, "KEY POINT", "Template creation is deterministic. Same input always produces same template. No randomness, no ML, no tuning.")

    # S6: VOXEL CODE
    s = light_slide(pres); pn = sn()
    add_badge(s, "CODE"); add_page_num(s, pn); add_title(s, "Voxel Downsampling: The Code"); add_subtitle(s, "Three lines of NumPy that reduce 55K to 22K points")
    add_body_text(s, "indices = np.floor(points / voxel_size).astype(np.int64)\n_, unique_idx = np.unique(indices, axis=0, return_index=True)\ndownsampled = points[unique_idx]\n\nResult: 55,000 \u2192 22,545 points\nEach voxel: 0.5mm\u00b3 cube\nPoints sharing the same cell collapse to one representative\n\nMeasurement requirement: \u00b10.001\" = \u00b10.025mm\nVoxel size: 0.5mm = 50x finer than requirement\nEngineering margin built into the discretization step", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "PRECISION", "0.5mm voxels give 50x finer resolution than our measurement requirement. That is deliberate engineering margin.")

    # S7: OUTLIER REMOVAL
    s = light_slide(pres); pn = sn()
    add_badge(s, "CODE"); add_page_num(s, pn); add_title(s, "Statistical Outlier Removal"); add_subtitle(s, "KD-tree k-NN distance filter removes scanner noise")
    add_body_text(s, "tree = cKDTree(points)\ndists, _ = tree.query(points, k=21)    # self + 20 neighbors\nmean_dists = np.mean(dists[:, 1:], axis=1)\nthreshold = global_mean + 2.5 * global_std\ninliers = mean_dists < threshold\n\nResult: 22,545 \u2192 22,064 points\n481 outliers removed (2.1%)\n\nEach point's 20-nearest-neighbor mean distance\ncompared to the global population statistics.\nPoints with unusually high mean distance are isolated noise.", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "WHY THIS MATTERS", "Scanner noise creates isolated points that would become phantom defects. This filter catches them before they corrupt analysis.")

    # S8: SECTION 02
    s = dark_slide(pres); sn()
    _dark_title(s, "02", "KD-Tree Mesh Network", "The spatial data structure that powers every query in the pipeline")

    # S9: WHAT IS KDTREE
    s = light_slide(pres); pn = sn()
    add_badge(s, "KD-TREE"); add_page_num(s, pn); add_title(s, "What Is a KD-Tree?"); add_subtitle(s, "Binary space partitioning for O(log n) spatial queries")
    add_body_text(s, "Binary tree that alternates splitting on X, Y, Z axes.\n\nBuild complexity:  O(n log n) \u2014 one-time cost\nQuery complexity:  O(log n) \u2014 per point\n\nBrute force on 22K points:\n  22,000 \u00d7 22,000 = 484,000,000 comparisons\n\nKD-tree on 22K points:\n  22,000 \u00d7 15 levels \u2248 330,000 comparisons\n\nSpeedup: ~1,500x\n\nOur implementation: scipy.spatial.cKDTree\nC-compiled. Multi-threaded via workers=-1.", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "PERFORMANCE", "Without KD-trees, Phase 3 takes 30 minutes. With them, 0.1 seconds. That is the difference between usable and unusable.")

    # S10: KDTREE ACROSS PIPELINE
    s = light_slide(pres); pn = sn()
    add_badge(s, "KD-TREE"); add_page_num(s, pn); add_title(s, "KD-Tree Across the Pipeline"); add_subtitle(s, "Used in 4 of 8 phases plus the defect library")
    add_card(s, "Phase 1: Outlier Removal", "k-NN distance to detect isolated scanner noise points.\nk=20 neighbors, compare to global mean+2.5\u03c3.", 0.5, 1.4, w=4.2, h=1.0)
    add_card(s, "Phase 2: ICP Registration", "Nearest-neighbor correspondence matching.\nIteratively finds closest CAD point for each scan point.", 5.2, 1.4, w=4.2, h=1.0)
    add_card(s, "Phase 3: Deviation Analysis", "Signed distance from every scan point to nearest CAD point.\nVectorized batch query: 22K points in one call.", 0.5, 2.6, w=4.2, h=1.0)
    add_card(s, "Defect Library: Adjacency", "query_ball_point for radius-based neighbor search.\nFinds all defects within 5mm of each other.", 5.2, 2.6, w=4.2, h=1.0)
    add_insight_bar(s, "BACKBONE", "The KD-tree is the computational backbone. Remove it and Phase 3 alone exceeds our 90-second budget.")

    # S11: SIGNED DISTANCE MATH
    s = light_slide(pres); pn = sn()
    add_badge(s, "MATH"); add_page_num(s, pn); add_title(s, "Signed Distance Computation"); add_subtitle(s, "The core mathematical operation for defect detection")
    add_body_text(s, "For each scan point p\u1d62:\n  1. Find nearest CAD point c\u2c7c = KD-tree.query(p\u1d62)\n  2. Get CAD surface normal n\u2c7c at c\u2c7c\n  3. signed_dist = (p\u1d62 \u2212 c\u2c7c) \u00b7 n\u2c7c\n\nConvention:\n  Negative = material removed (nick, gouge, crack)\n  Positive = material added (buildup, deposit)\n\nVectorized implementation:\n  tree = cKDTree(cad_pts)\n  _, indices = tree.query(scan_pts, k=1, workers=-1)\n  diff = scan_pts - cad_pts[indices]\n  signed = np.sum(diff * cad_normals[indices], axis=1)", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "VECTORIZED", "No Python loops. 22,064 signed distances computed in one NumPy dot product operation.")

    # S12: DEVIATION RESULTS
    s = light_slide(pres); pn = sn()
    add_badge(s, "RESULTS"); add_page_num(s, pn); add_title(s, "Deviation Analysis Results"); add_subtitle(s, "Output from our first pipeline run on synthetic IBR data")
    add_stat_box(s, "22,064", "Points\nanalyzed", 0.3, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "\u00b10.14", "Deviation range\n(mm)", 2.65, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "800", "Below threshold\n(3.6%)", 5.0, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "4", "Defect clusters\nafter DBSCAN", 7.35, 1.4, w=2.15, h=1.3)
    add_body_text(s, "Mean deviation: -0.0008mm (near zero = good alignment)\nStd deviation: 0.0132mm\nThreshold: -0.010mm = 2x scanner precision floor", x=0.5, y=3.0, h=0.8, size=11, color=BODY)
    add_insight_bar(s, "VALIDATION", "3.6% candidate rate on synthetic data with 4 injected defects. Threshold performing as designed.")

    # S13: SECTION 03
    s = dark_slide(pres); sn()
    _dark_title(s, "03", "Zone Template System", "Mapping every point on the blade to P&W's 13-zone classification")

    # S14: TWO-LAYER ZONES
    s = light_slide(pres); pn = sn()
    add_badge(s, "ZONES"); add_page_num(s, pn); add_title(s, "Two-Layer Zone Boundary Design"); add_subtitle(s, "Percentages for portability, millimeters for precision")
    add_card(s, "Layer 1: Percentages (Stored)", "A1: 0\u20133.83% of blade span\nA2: 3.83\u201362.88% of span\nA3: 62.88\u201399.77% of span\nPortable across all 9 F135 stages", 0.5, 1.4, w=4.2, h=1.4)
    add_card(s, "Layer 2: Millimeters (Runtime)", "resolve_zone(pct, measured_span_mm)\nUses ACTUAL measured blade span from ICP\nHandles manufacturing variation\nWorn blade gets correct boundaries", 5.2, 1.4, w=4.2, h=1.4)
    add_body_text(s, "Why measured span, not nominal?\nICP alignment gives actual as-manufactured geometry.\nA blade worn 2mm at the tip has different effective span.", x=0.5, y=3.1, h=0.8, size=11, color=BODY)
    add_insight_bar(s, "PORTABILITY", "One config file works across all 9 F135 stages. Zero per-stage calibration needed.")

    # S15: ZONE CLASSIFICATION LOGIC
    s = light_slide(pres); pn = sn()
    add_badge(s, "LOGIC"); add_page_num(s, pn); add_title(s, "Zone Classification Decision Tree"); add_subtitle(s, "How each defect gets mapped to one or more P&W zones")
    add_body_text(s, "if edge and nearest_edge == 'LE':\n    match A1/A2/A3 by span percentage\nelif edge and nearest_edge == 'TE':\n    match B1/B2 by span percentage\nelse:  # surface\n    default C1 (convex) or C2 (concave)\n\nALWAYS check tip zones regardless of classification:\n  tip_zones = match_tip_zones(height_pct)\n  # A defect at 95% span on LE gets both A3 AND E2\n\nMulti-zone resolution:\n  most_restrictive = min(limits_list,\n      key=lambda l: (l.max_depth_in, l.max_length_in))\n\nConfirmed by P&W: always most restrictive.", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "TIP ZONE CHECK", "Tip zone check is independent of edge/surface. A defect near the tip always gets tightest applicable limits.")

    # S16: COMPLIANCE AND LOGIC
    s = light_slide(pres); pn = sn()
    add_badge(s, "COMPLIANCE"); add_page_num(s, pn); add_title(s, "Compliance Check: AND Logic, Three Tiers"); add_subtitle(s, "Every condition must pass simultaneously")
    add_card(s, "SERVICEABLE", "depth \u2264 max_depth AND\nlength \u2264 max_length", 0.5, 1.4, w=2.8, h=0.9)
    add_card(s, "BLEND", "depth \u2264 1.5\u00d7 max_depth AND\nlength \u2264 1.5\u00d7 max_length", 3.6, 1.4, w=2.8, h=0.9)
    add_card(s, "REPLACE", "Exceeds blend limits\nBlade replacement required", 6.7, 1.4, w=2.8, h=0.9)
    add_body_text(s, "Special modifiers:\n  Cracks: 0.5x normal limits (half tolerance)\n  FOD: 0.8x normal limits\n  Unknown zone: default to most restrictive\n\nAND logic = a defect 0.001\" over on length but\nperfect on depth still gets flagged. No partial passes.", x=0.5, y=2.6, h=1.4, size=11, color=BODY)
    add_insight_bar(s, "AEROSPACE SAFETY", "No partial passes. Every dimension must be within limits simultaneously. This is non-negotiable.")

    # S17: SECTION 04
    s = dark_slide(pres); sn()
    _dark_title(s, "04", "Defect Library: Hybrid Storage", "Grid + Metric dual-index for cross-dimensional analysis")

    # S18: DEFECT LIBRARY ARCH
    s = light_slide(pres); pn = sn()
    add_badge(s, "LIBRARY"); add_page_num(s, pn); add_title(s, "DefectLibrary Architecture"); add_subtitle(s, "Four-dimensional query system with optimized indices")
    add_body_text(s, "Primary store:  dict[defect_id \u2192 defect_dict]\n\nSecondary indices:\n  by_foil: foil_number \u2192 [defect_ids]    (instant lookup)\n  by_zone: zone_id \u2192 [defect_ids]        (instant lookup)\n\nSpatial index:\n  cKDTree on defect centroids\n  Rebuilt lazily via _tree_dirty flag\n\n4 Dimensions of Analysis:\n  D1: WHERE is the defect  (Grid O(1) or Metric O(log n))\n  D2: WHICH zone           (percentage-based lookup)\n  D3: CROSS-ZONE           (aggregate via secondary index)\n  D4: CROSS-FOIL           (RotorDefectAnalysis class)", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "QUERY OPTIMIZATION", "Each dimension has its own optimized query path. D1-D3 run per-blade in parallel. D4 runs once after all blades finish.")

    # S19: ADJACENCY CODE
    s = light_slide(pres); pn = sn()
    add_badge(s, "CODE"); add_page_num(s, pn); add_title(s, "Adjacency Detection: The Code"); add_subtitle(s, "KD-tree radius query for cross-defect proximity analysis")
    add_body_text(s, "def find_adjacent_defects(self, defect_id, threshold_mm=5.0):\n    center = self._centroids[defect_id]\n    indices = self._kdtree.query_ball_point(\n        center, r=threshold_mm)\n    return [defects[tree_ids[i]] for i in indices\n            if tree_ids[i] != defect_id]\n\nWhy adjacency matters:\n  Two small defects 3mm apart may each pass\n  individual zone limits, but together they create\n  a structural weakness that needs flagging.\n\n  query_ball_point: O(log n) radius search.\n  No brute-force distance matrix needed.", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "PROXIMITY", "Cross-defect proximity is a safety-critical check. Two passing defects can together fail if they are close enough.")

    # S20: SECTION 05
    s = dark_slide(pres); sn()
    _dark_title(s, "05", "Results Walkthrough", "Tracing defects through the full pipeline end to end")

    # S21: SINGLE DEFECT TRACE
    s = light_slide(pres); pn = sn()
    add_badge(s, "TRACE"); add_page_num(s, pn); add_title(s, "Defect F001_D002: Full 8-Phase Trace"); add_subtitle(s, "Following one defect from raw scan to REPLACE disposition")
    add_body_text(s, "Phase 1: Part of 55,000 raw pts \u2192 survives voxel + outlier\nPhase 2: Aligned with 0.016mm RMSE \u2192 correct position\nPhase 3: 345 points have signed distance -0.048 to -0.122mm\n         All below -0.010mm threshold \u2192 flagged as candidates\nPhase 5: DBSCAN groups 345 nearby pts into one cluster\n         eps=1.0mm, min_samples=3 \u2192 one contiguous defect\nPhase 6: OBB measurement:\n         Length = 0.5960\" | Width = 0.4683\" | Depth = 0.0048\"\nPhase 7: Zones C1 + E2\n         C1 limit: depth\u22640.010\", length\u22640.080\"\n         Length 0.596\" >> 0.080\" \u2192 REPLACE", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "TRACEABILITY", "Remove any single phase and this defect gets missed or gets wrong disposition. Every phase contributes.")

    # S22: ALL 4 DEFECTS TABLE
    s = light_slide(pres); pn = sn()
    add_badge(s, "RESULTS"); add_page_num(s, pn); add_title(s, "All 4 Detected Defects"); add_subtitle(s, "Complete results from the initial pipeline run")
    _add_table(s, ["ID", "Type", "Zones", "Depth (in)", "Length (in)", "Disposition"], [
        ["F001_D001", "Surface", "C1, E2", "0.0011", "0.0757", "SERVICEABLE"],
        ["F001_D002", "Surface", "C1, E2", "0.0048", "0.5960", "REPLACE"],
        ["F001_D003", "Surface", "C1, E2", "0.0056", "0.6105", "REPLACE"],
        ["F001_D004", "Surface", "C1, E2", "0.0006", "0.0364", "SERVICEABLE"],
    ])
    add_insight_bar(s, "MATCH", "D002/D003 exceeded C1 max_length 0.080\" by 7x. System correctly identified them as REPLACE. Matches expected outcomes.")

    # S23: PCA vs OBB
    s = light_slide(pres); pn = sn()
    add_badge(s, "MEASUREMENT"); add_page_num(s, pn); add_title(s, "PCA vs OBB: Why It Matters"); add_subtitle(s, "Domain-validated measurement distinction from P&W")
    add_card(s, "Edge Defects \u2192 PCA", "Within 0.080\" of LE/TE\nPrincipal component aligns with edge tangent\nMeasures LONGITUDINALLY along the curve\nEdge nicks elongate along the edge", 0.5, 1.4, w=4.2, h=1.5)
    add_card(s, "Surface Defects \u2192 OBB", "Away from edges (>0.080\")\nOriented Bounding Box via convex hull\nMeasures TRANSVERSELY across the surface\nSurface dents are more isotropic", 5.2, 1.4, w=4.2, h=1.5)
    add_body_text(s, "Using the wrong method gives wrong measurements.\nThis distinction came from Bryan at P&W, not from the math.", x=0.5, y=3.15, h=0.6, size=11, color=BODY)
    add_insight_bar(s, "DOMAIN KNOWLEDGE", "PCA vs OBB is a domain decision you cannot derive from algorithms alone. Required direct input from P&W engineers.")

    # S24: SECTION 06
    s = dark_slide(pres); sn()
    _dark_title(s, "06", "FastAPI Dashboard", "Live interactive web application code walkthrough")

    # S25: FASTAPI STACK
    s = light_slide(pres); pn = sn()
    add_badge(s, "FASTAPI"); add_page_num(s, pn); add_title(s, "FastAPI Architecture & Key Code"); add_subtitle(s, "Async Python serving Plotly 3D visualizations over REST")
    add_body_text(s, "Stack: FastAPI (async) \u2192 Jinja2 templates \u2192 Plotly.js (WebGL)\n\nKey endpoint — deviation heatmap data:\n\n@app.get('/api/pointcloud/deviations')\nasync def api_deviations():\n    scan_pcd = read_point_cloud(scan_path)\n    cad_pcd = read_point_cloud(cad_path)\n    tree = cKDTree(cad_pts)\n    _, idx = tree.query(scan_pts, k=1, workers=-1)\n    diff = scan_pts - cad_pts[idx]\n    signed = np.sum(diff * normals[idx], axis=1)\n    return {'points': pts.tolist(),\n            'deviations': signed.tolist()}", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "SAME MATH", "Same KD-tree computation from Phase 3 exposed as a REST API. 30,000 signed distances computed per API call.")

    # S26: DASHBOARD ELEMENTS
    s = light_slide(pres); pn = sn()
    add_badge(s, "DASHBOARD"); add_page_num(s, pn); add_title(s, "Dashboard: Visual Elements"); add_subtitle(s, "Every chart reads from actual pipeline output — nothing mocked")
    add_card(s, "Stats Row", "5 cards: Part#, Defects, Serviceable, Blend, Replace", 0.5, 1.4, w=4.2, h=0.7)
    add_card(s, "3D Deviation Heatmap", "Plotly scatter3d, 30K points, blue\u2192green\u2192red colorscale", 5.2, 1.4, w=4.2, h=0.7)
    add_card(s, "Disposition Donut", "Pie chart: green/amber/red segments with counts", 0.5, 2.25, w=4.2, h=0.7)
    add_card(s, "Zone Distribution", "Bar chart showing defect count per P&W zone", 5.2, 2.25, w=4.2, h=0.7)
    add_card(s, "Defect Table", "All defects with measurements and color-coded badges", 0.5, 3.1, w=4.2, h=0.7)
    add_card(s, "Depth Comparison", "Bar chart: each defect depth vs zone limit threshold", 5.2, 3.1, w=4.2, h=0.7)
    add_insight_bar(s, "REAL DATA", "Re-run the pipeline, refresh the browser — dashboard updates instantly. No cache, no stale data.")

    # S27: 3D VIEWER CODE
    s = light_slide(pres); pn = sn()
    add_badge(s, "3D VIEWER"); add_page_num(s, pn); add_title(s, "3D Viewer: Plotly.js Rendering"); add_subtitle(s, "Client-side WebGL handles 50K points smoothly in the browser")
    add_body_text(s, "4 View Modes:\n  1. Deviation Heatmap \u2014 color by signed distance\n  2. Scan Point Cloud \u2014 raw scan in green\n  3. CAD Reference \u2014 blue reference model\n  4. Defect Clusters \u2014 colored by defect ID\n\nPlotly.js rendering:\n\nPlotly.newPlot('viewer', [{\n  x, y, z,\n  mode: 'markers', type: 'scatter3d',\n  marker: {\n    size: 2, color: deviations,\n    colorscale: [['#3b82f6'],['#22c55e'],['#ef4444']]\n  }\n}]);\n\nControls: point size slider, max points, reset camera", x=0.5, y=1.4, h=2.6, size=11)
    add_insight_bar(s, "BROWSER GPU", "No server-side rendering. Plotly uses WebGL — the browser GPU does all the 3D work. Server only sends point data.")

    # S28: API ENDPOINTS TABLE
    s = light_slide(pres); pn = sn()
    add_badge(s, "API"); add_page_num(s, pn); add_title(s, "12 REST API Endpoints"); add_subtitle(s, "Integration surface for Hitachi systems and external tools")
    _add_table(s, ["Endpoint", "Method", "Description"], [
        ["/api/report/latest", "GET", "Latest inspection report JSON"],
        ["/api/reports", "GET", "List all reports with metadata"],
        ["/api/download/{file}", "GET", "Download JSON or Excel file"],
        ["/api/pointcloud/scan", "GET", "Scan cloud (40K points)"],
        ["/api/pointcloud/deviations", "GET", "Deviation-colored cloud"],
        ["/api/config", "GET", "Pipeline YAML configuration"],
        ["/api/rotor-configs", "GET", "F135 rotor definitions"],
        ["/api/pipeline/run", "POST", "Trigger pipeline execution"],
        ["/api/pipeline/status", "GET", "Pipeline run progress"],
    ], h=2.8)
    add_insight_bar(s, "INTEGRATION", "Any system can POST to /api/pipeline/run and poll /api/pipeline/status. Ready for Hitachi system integration.")

    # S29: DISCUSSION QUESTIONS
    s = light_slide(pres); pn = sn()
    add_badge(s, "DISCUSSION"); add_page_num(s, pn); add_title(s, "Discussion Questions for the Team"); add_subtitle(s, "Technical decisions that need group input before Sprint 4")
    for i, (t, d) in enumerate([("Voxel Resolution", "Should we go finer than 0.5mm for production scanner data?"),
        ("KD-tree vs R-tree", "Would R-tree be better for zone boundary polygon queries?"),
        ("LE/TE Curve Extraction", "How do we extract edge curves from CAD mesh automatically?"),
        ("Multi-Zone Handling", "Defects spanning 3+ zones \u2014 flag for manual review?"),
        ("Dashboard Integration", "What other Hitachi systems need API access?"),
        ("Real Scan Validation", "When can we validate against actual P&W scanner output?")]):
        add_card(s, f"{i+1}. {t}", d, 0.5, 1.35 + i*0.47, w=9, h=0.42)
    add_insight_bar(s, "BLOCKERS", "Items 3 and 6 are blocking. I need answers on LE/TE extraction and real scan data before Sprint 4 planning.")

    # S30: CLOSING
    s = dark_slide(pres); sn()
    t = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8))
    r = t.text_frame.paragraphs[0].add_run()
    r.text = "The Template Is Built.\nThe Network Is Running."; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE; r.font.name = "Georgia"
    t2 = s.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(8.4), Inches(0.4))
    r2 = t2.text_frame.paragraphs[0].add_run()
    r2.text = "Grid discretization + KD-tree mesh = hybrid precision at production speed"; r2.font.size = Pt(14); r2.font.color.rgb = PRIMARY_LIGHT; r2.font.name = "Arial"
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.7), Inches(3), Inches(0.04))
    ln.fill.solid(); ln.fill.fore_color.rgb = PRIMARY_LIGHT; ln.line.fill.background()
    for i, txt in enumerate(["Supreeth Mysore  |  supreeth.mysore@hitachids.com", "Hitachi Digital Services  |  P&W F135 Program", "github.com/drmysore/open3d-ibr-defect-extraction"]):
        _closing_line(s, txt, 3.0 + i*0.35)

    out = os.path.join("output", "IBR_Technical_DeepDive.pptx")
    os.makedirs("output", exist_ok=True)
    pres.save(out)
    print(f"Saved: {out} ({n[0]} slides)")


if __name__ == "__main__":
    build()
