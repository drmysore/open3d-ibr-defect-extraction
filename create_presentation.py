"""Generate IBR Defect Extraction System presentation.

Author: Supreeth Mysore, Lead GPU Architect
Program: F135 IBR Automated Defect Extraction
Organization: Hitachi Digital Services
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
import os

# ── DESIGN TOKENS ──
PRIMARY = RGBColor(0x67, 0x09, 0x45)
PRIMARY_LIGHT = RGBColor(0x8B, 0x2D, 0x5E)
PRIMARY_DARK = RGBColor(0x4A, 0x06, 0x32)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TINT = RGBColor(0xF9, 0xF5, 0xF7)
DARK = RGBColor(0x33, 0x33, 0x33)
BODY = RGBColor(0x66, 0x66, 0x66)
MUTED = RGBColor(0x99, 0x99, 0x99)
BORDER = RGBColor(0xE0, 0xE0, 0xE0)
GREEN_ACCENT = RGBColor(0x22, 0x8B, 0x22)
AMBER_ACCENT = RGBColor(0xD4, 0x8B, 0x00)
RED_ACCENT = RGBColor(0xC0, 0x39, 0x2B)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)
TOTAL_SLIDES = 24


def add_badge(slide, text, x=0.5, y=0.22):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.8), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text.upper()
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"


def add_page_num(slide, num):
    txBox = slide.shapes.add_textbox(Inches(8.3), Inches(0.25), Inches(1.2), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{num} / {TOTAL_SLIDES}"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED
    run.font.name = "Arial"


def add_title(slide, text, y=0.7, size=30):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    run.font.name = "Georgia"


def add_subtitle(slide, text, y=1.08):
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(y), Inches(9), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(13)
    run.font.color.rgb = BODY
    run.font.name = "Arial"


def add_insight_bar(slide, label, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(4.35), SLIDE_W, Inches(1.275))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(4.48), Inches(8.8), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = label.upper()
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Arial"

    txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(4.78), Inches(8.8), Inches(0.6))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = text
    run2.font.size = Pt(11)
    run2.font.color.rgb = WHITE
    run2.font.name = "Arial"


def add_body_text(slide, text, x=0.5, y=1.4, w=9, h=2.5, size=12, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.name = "Arial"


def add_stat_box(slide, value, label, x, y, w=2.5, h=1.5):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = TINT
    shape.line.color.rgb = BORDER

    txBox = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.15), Inches(w - 0.2), Inches(0.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(value)
    run.font.size = Pt(42)
    run.font.bold = True
    run.font.color.rgb = PRIMARY
    run.font.name = "Georgia"

    txBox2 = slide.shapes.add_textbox(Inches(x + 0.1), Inches(y + 0.9), Inches(w - 0.2), Inches(0.4))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(11)
    run2.font.color.rgb = BODY
    run2.font.name = "Arial"


def add_card(slide, title, body, x, y, w=4.2, h=1.2):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = BORDER

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY
    accent.line.fill.background()

    txBox = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.1), Inches(w - 0.35), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(12)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = "Arial"

    txBox2 = slide.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.4), Inches(w - 0.35), Inches(h - 0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = body
    run2.font.size = Pt(10)
    run2.font.color.rgb = BODY
    run2.font.name = "Arial"


def add_process_step(slide, num, title, desc, x, y):
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(num)
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Georgia"

    txBox = slide.shapes.add_textbox(Inches(x - 0.2), Inches(y + 0.55), Inches(0.85), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.color.rgb = DARK
    run.font.name = "Arial"

    txBox2 = slide.shapes.add_textbox(Inches(x - 0.3), Inches(y + 0.85), Inches(1.05), Inches(0.5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(8)
    run2.font.color.rgb = BODY
    run2.font.name = "Arial"


def dark_slide(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = PRIMARY_DARK
    return slide


def light_slide(pres):
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide


def build_presentation():
    pres = Presentation()
    pres.slide_width = SLIDE_W
    pres.slide_height = SLIDE_H
    n = [0]

    def sn():
        n[0] += 1
        return n[0]

    # ───── SLIDE 1: TITLE (DARK) ─────
    s = dark_slide(pres)
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "IBR Automated Defect Extraction System"
    run.font.size = Pt(34)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.4), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "Open3D 8-Phase Pipeline  |  Sprint 3 Initial Results"
    run2.font.size = Pt(16)
    run2.font.color.rgb = PRIMARY_LIGHT
    run2.font.name = "Arial"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.5), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    for i, txt in enumerate([
        "Supreeth Mysore  |  Lead GPU Architect",
        "Hitachi Digital Services  |  Pratt & Whitney F135 Program",
        "March 2026  |  Sprint 3 Review"
    ]):
        txBox = s.shapes.add_textbox(Inches(0.8), Inches(2.8 + i * 0.35), Inches(8.4), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        run.font.name = "Arial"

    sn()

    # ───── SLIDE 2: AGENDA ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "OVERVIEW")
    add_page_num(s, pn)
    add_title(s, "Agenda")
    add_subtitle(s, "30-minute walkthrough of our Sprint 3 progress and initial pipeline results")

    items = [
        ("1. Problem & Motivation", "Why we are building this; the manual inspection bottleneck"),
        ("2. Architecture & Pipeline", "The 8-phase Open3D processing pipeline end-to-end"),
        ("3. Initial Results", "First run metrics, defect detection, zone classification"),
        ("4. Interactive Dashboard", "Live FastAPI web application with 3D visualization"),
        ("5. Next Steps & Sprint 4", "What we still need, open questions for P&W"),
    ]
    for i, (title, desc) in enumerate(items):
        add_card(s, title, desc, 0.5, 1.4 + i * 0.55, w=9, h=0.48)

    add_insight_bar(s, "MEETING GOAL", "Walk through the working pipeline, demonstrate initial results, and align on Sprint 4 priorities with the team.")

    # ───── SLIDE 3: PROBLEM (SECTION DIVIDER) ─────
    s = dark_slide(pres)
    pn = sn()
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "01"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "The Problem We Are Solving"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Georgia"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = "Manual blade inspection is a throughput bottleneck in F135 MRO operations"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    run3.font.name = "Arial"

    # ───── SLIDE 4: MANUAL VS AUTOMATED ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "PROBLEM")
    add_page_num(s, pn)
    add_title(s, "Manual Inspection: The Bottleneck")
    add_subtitle(s, "Current state versus what we are building")

    add_stat_box(s, "5-12", "Minutes per blade\n(manual)", 0.5, 1.4, w=2.8, h=1.3)
    add_stat_box(s, "11 hrs", "Stage 3 IBR\n(55 blades)", 3.6, 1.4, w=2.8, h=1.3)
    add_stat_box(s, "22 hrs", "Stage 9 IBR\n(110 blades)", 6.7, 1.4, w=2.8, h=1.3)

    add_body_text(s, "Manual process: visual exam under magnification, caliper measurement,\ncomparison against reparable limits, written documentation per blade.\nHuman variability: \u00b10.003\" typical vs \u00b10.001\" required.", x=0.5, y=3.0, h=1.0, size=11, color=BODY)

    add_insight_bar(s, "THE BOTTLENECK", "A single Stage 9 IBR takes a full technician day. Our system targets under 90 seconds for the same work at higher precision.")

    # ───── SLIDE 5: TARGET STATE ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "SOLUTION")
    add_page_num(s, pn)
    add_title(s, "Automated Inspection: Our Target")
    add_subtitle(s, "Key performance targets for the production system")

    stats = [("<90s", "Per IBR\nprocessing time"), ("\u00b10.001\"", "Measurement\nprecision"), ("100%", "Defect\ndetection rate"), ("<2%", "False positive\nrate")]
    for i, (val, lbl) in enumerate(stats):
        add_stat_box(s, val, lbl, 0.3 + i * 2.35, 1.4, w=2.15, h=1.3)

    add_body_text(s, "Replaces manual inspection with 3D LiDAR scan processing.\nProcesses entire IBR in one pass. Generates P&W-format compliance reports.\nCost: $0.08-$0.34 per IBR inspection on AWS.", x=0.5, y=3.0, h=1.0, size=11, color=BODY)
    add_insight_bar(s, "COST IMPACT", "$0.18 per Stage 3 IBR versus hours of technician time. Annual projection at 720 IBRs: under $250 total compute cost.")

    # ───── SLIDE 6: SECTION DIVIDER - ARCHITECTURE ─────
    s = dark_slide(pres)
    pn = sn()
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "02"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    tf2 = txBox2.text_frame
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "System Architecture"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Georgia"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    run3 = tf3.paragraphs[0].add_run()
    run3.text = "8-phase Open3D pipeline from raw PLY scan to P&W compliance report"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    run3.font.name = "Arial"

    # ───── SLIDE 7: 8-PHASE PIPELINE FLOW ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "PIPELINE")
    add_page_num(s, pn)
    add_title(s, "8-Phase Processing Pipeline")
    add_subtitle(s, "End-to-end from PLY scan to compliance report")

    phases = [
        ("1", "Data\nPrep", "Load, clean,\ndownsample"),
        ("2", "Register", "RANSAC +\nICP align"),
        ("3", "Deviate", "KD-tree\nsigned dist"),
        ("4", "Segment", "Angular\nDBSCAN"),
        ("5", "Cluster", "Spatial\nDBSCAN"),
        ("6", "Measure", "PCA / OBB\ndimensions"),
        ("7", "Classify", "13-zone +\ncompliance"),
        ("8", "Report", "Excel +\nJSON output"),
    ]
    for i, (num, title, desc) in enumerate(phases):
        add_process_step(s, num, title, desc, 0.5 + i * 1.15, 1.45)

    add_insight_bar(s, "ARCHITECTURE", "Each phase is independently configurable via YAML. No magic numbers in code. All 91 compliance rules externalized.")

    # ───── SLIDE 8: REGISTRATION DETAIL ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "PHASE 2")
    add_page_num(s, pn)
    add_title(s, "Registration: Scan-to-CAD Alignment")
    add_subtitle(s, "Two-stage approach: coarse global + fine local alignment")

    add_card(s, "Stage 1: RANSAC (Coarse)", "Feature-based global alignment using FPFH descriptors.\n4M iterations, 0.999 confidence. Handles wildly off scans.", 0.5, 1.4, w=4.2, h=1.2)
    add_card(s, "Stage 2: ICP (Fine)", "Iterative Closest Point refines to sub-micron precision.\nPoint-to-Plane estimation, 200 max iterations.", 5.2, 1.4, w=4.2, h=1.2)
    add_card(s, "Validation", "RMSE must be < 0.05mm. If alignment fails, inspection is flagged.\nOur initial run achieved 0.016mm RMSE.", 0.5, 2.8, w=4.2, h=1.0)
    add_card(s, "Why This Matters", "Misaligned scan = every point looks like a defect.\nRegistration quality gates everything downstream.", 5.2, 2.8, w=4.2, h=1.0)

    add_insight_bar(s, "INITIAL RESULT", "Registration RMSE of 0.016mm on synthetic data validates the two-stage approach. Well within our 0.05mm acceptance threshold.")

    # ───── SLIDE 9: DEVIATION ANALYSIS ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "PHASE 3")
    add_page_num(s, pn)
    add_title(s, "Deviation Analysis: Finding Defects")
    add_subtitle(s, "Vectorized KD-tree signed distance computation")

    add_body_text(s, "For every scan point, compute signed distance to nearest CAD surface:\n\n  \u2022  Negative deviation = material removed (nick, gouge, crack)\n  \u2022  Positive deviation = material added (buildup, deposit)\n  \u2022  Threshold: -0.010mm selected for 100% detection, 1.8% false positive rate\n\nSciPy cKDTree provides 10x speedup over loop-based approach.\nHandles 20M points in under 50 seconds.", x=0.5, y=1.4, h=2.5, size=12)

    add_insight_bar(s, "THRESHOLD RATIONALE", "-0.010mm gives 100% detection with 1.8% FP rate. Provides 2x margin above scanner precision floor of \u00b10.005mm.")

    # ───── SLIDE 10: THRESHOLD ANALYSIS TABLE ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "SENSITIVITY")
    add_page_num(s, pn)
    add_title(s, "Sensitivity Analysis: Threshold Selection")
    add_subtitle(s, "Systematic evaluation of detection vs false positive tradeoff")

    rows = [
        ["-0.005mm", "100%", "12.4%", "Too sensitive"],
        ["-0.008mm", "100%", "4.1%", "Acceptable"],
        ["-0.010mm", "100%", "1.8%", "SELECTED"],
        ["-0.012mm", "99.2%", "0.9%", "Misses defects"],
        ["-0.015mm", "96.1%", "0.3%", "Unacceptable"],
    ]
    headers = ["Threshold", "Detection", "False Positive", "Assessment"]

    table_shape = s.shapes.add_table(len(rows) + 1, 4, Inches(0.5), Inches(1.4), Inches(9), Inches(2.4))
    table = table_shape.table
    table.columns[0].width = Inches(2)
    table.columns[1].width = Inches(2)
    table.columns[2].width = Inches(2.5)
    table.columns[3].width = Inches(2.5)

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TINT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = "Arial"
                    r.font.color.rgb = DARK
                    if ri == 2:
                        r.font.bold = True
                p.alignment = PP_ALIGN.CENTER

    add_insight_bar(s, "DECISION", "We went with -0.010mm. Zero missed defects is non-negotiable in aerospace. 1.8% FP rate gets filtered by DBSCAN clustering downstream.")

    # ───── SLIDE 11: 13-ZONE SYSTEM ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "ZONES")
    add_page_num(s, pn)
    add_title(s, "P&W 13-Zone Classification System")
    add_subtitle(s, "Zone-specific reparable limits drive disposition decisions")

    zone_data = [
        ("LE_TIP / TE_TIP", "CRITICAL", "0.005\"", "0.030\""),
        ("A1 / B1 (LE/TE Critical)", "CRITICAL", "0.003\"", "0.020\""),
        ("A2, A3, B2 (LE/TE Mid/Outer)", "HIGH", "0.005-0.008\"", "0.040-0.060\""),
        ("C1, C2 (Convex/Concave)", "STANDARD", "0.010\"", "0.080\""),
        ("D1, D2 (Platform)", "STANDARD", "0.015\"", "0.100\""),
        ("E1, E2 (Root/Tip)", "STANDARD", "0.008-0.010\"", "0.060-0.080\""),
    ]
    headers2 = ["Zone Group", "Severity", "Max Depth", "Max Length"]
    table_shape2 = s.shapes.add_table(len(zone_data) + 1, 4, Inches(0.5), Inches(1.4), Inches(9), Inches(2.6))
    table2 = table_shape2.table
    table2.columns[0].width = Inches(3.5)
    table2.columns[1].width = Inches(1.5)
    table2.columns[2].width = Inches(2)
    table2.columns[3].width = Inches(2)

    for i, h in enumerate(headers2):
        cell = table2.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(zone_data):
        for ci, val in enumerate(row):
            cell = table2.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TINT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = "Arial"
                    r.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    add_insight_bar(s, "MULTI-ZONE RULE", "When a defect spans two zones, we apply the most restrictive limit. Confirmed directly by P&W. Non-negotiable for aerospace safety.")

    # ───── SLIDE 12: COMPLIANCE ENGINE ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "COMPLIANCE")
    add_page_num(s, pn)
    add_title(s, "Compliance Engine: 91 Rules")
    add_subtitle(s, "13 zones x 7 defect types, three disposition tiers")

    add_card(s, "SERVICEABLE", "Depth and length within zone limits.\nReturn to service with no action needed.", 0.5, 1.4, w=2.8, h=1.1)
    add_card(s, "BLEND", "Within 1.5x zone limits.\nBlend/repair the defect and recheck.", 3.6, 1.4, w=2.8, h=1.1)
    add_card(s, "REPLACE", "Exceeds blend limits.\nBlade or rotor replacement required.", 6.7, 1.4, w=2.8, h=1.1)

    add_body_text(s, "AND logic: all conditions must pass simultaneously.\nCracks get 0.5x normal limits. FOD gets 0.8x.\n7 defect types: nick, dent, crack, FOD, erosion, scratch, gouge.", x=0.5, y=2.8, h=1.2, size=11, color=BODY)

    add_insight_bar(s, "SAFETY FIRST", "Unknown zones default to the most restrictive limits. We would rather flag a false positive than miss a real crack in a turbine blade.")

    # ───── SLIDE 13: SECTION DIVIDER - RESULTS ─────
    s = dark_slide(pres)
    pn = sn()
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "03"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    tf2 = txBox2.text_frame
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "Initial Pipeline Results"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Georgia"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    run3 = tf3.paragraphs[0].add_run()
    run3.text = "First end-to-end run on synthetic IBR data with controlled defects"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    run3.font.name = "Arial"

    # ───── SLIDE 14: FIRST RUN METRICS ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "RESULTS")
    add_page_num(s, pn)
    add_title(s, "First Pipeline Run: Key Metrics")
    add_subtitle(s, "Synthetic 5-blade IBR with controlled nick, dent, and gouge defects")

    add_stat_box(s, "14.5s", "Total\nprocessing time", 0.3, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "0.016mm", "Registration\nRMSE", 2.65, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "4", "Defects\ndetected", 5.0, 1.4, w=2.15, h=1.3)
    add_stat_box(s, "3.6%", "Defect\ncandidates", 7.35, 1.4, w=2.15, h=1.3)

    add_body_text(s, "Input: 55,000 point synthetic IBR with 5 blades.\nDownsampled to 22,064 points. 800 defect candidates identified.\n4 discrete defect clusters confirmed after DBSCAN filtering.\nDisposition: 2 SERVICEABLE, 2 REPLACE.", x=0.5, y=3.0, h=1.0, size=11, color=BODY)

    add_insight_bar(s, "VALIDATION", "14.5 seconds on a single CPU with no GPU. Well within the 90-second target. Production AWS GPU instances will be significantly faster.")

    # ───── SLIDE 15: DEFECT TABLE ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "DEFECTS")
    add_page_num(s, pn)
    add_title(s, "Detected Defect Details")
    add_subtitle(s, "All four defects from the initial pipeline run")

    defect_rows = [
        ["F001_D001", "Surface", "C1, E2", "0.0011\"", "0.0757\"", "SERVICEABLE"],
        ["F001_D002", "Surface", "C1, E2", "0.0048\"", "0.5960\"", "REPLACE"],
        ["F001_D003", "Surface", "C1, E2", "0.0056\"", "0.6105\"", "REPLACE"],
        ["F001_D004", "Surface", "C1, E2", "0.0006\"", "0.0364\"", "SERVICEABLE"],
    ]
    d_headers = ["Defect ID", "Type", "Zones", "Depth", "Length", "Disposition"]
    dt = s.shapes.add_table(5, 6, Inches(0.5), Inches(1.4), Inches(9), Inches(2.2))
    dtable = dt.table
    for ci, h in enumerate(d_headers):
        cell = dtable.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(10)
                r.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(defect_rows):
        for ci, val in enumerate(row):
            cell = dtable.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TINT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)
                    r.font.name = "Arial"
                    r.font.color.rgb = DARK
                    if ci == 5 and val == "REPLACE":
                        r.font.bold = True
                p.alignment = PP_ALIGN.CENTER

    add_insight_bar(s, "DISPOSITION", "Two defects exceeded zone C1 limits (max length 0.080\"). The system correctly identified them as REPLACE candidates.")

    # ───── SLIDE 16: MEASUREMENT APPROACH ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "MEASUREMENT")
    add_page_num(s, pn)
    add_title(s, "Measurement: PCA vs OBB")
    add_subtitle(s, "Domain-validated distinction confirmed by Bryan (P&W)")

    add_card(s, "Edge Defects: PCA", "Within 0.080\" of LE/TE. Measure longitudinally\nalong the edge curve. First principal component\naligns with the natural edge direction.", 0.5, 1.4, w=4.2, h=1.4)
    add_card(s, "Surface Defects: OBB", "Away from edges. Measure transversely using\nOriented Bounding Box. Surface defects are\nmore isotropic than edge defects.", 5.2, 1.4, w=4.2, h=1.4)

    add_body_text(s, "Using the wrong method produces wrong measurements.\nThis was a critical domain decision, not an engineering shortcut.\nEdge threshold: 0.080\" (2.032mm) from LE or TE curve.", x=0.5, y=3.1, h=0.9, size=11, color=BODY)

    add_insight_bar(s, "DOMAIN KNOWLEDGE", "This PCA vs OBB distinction is one of those decisions that you cannot figure out from the math alone. Took direct domain input to get right.")

    # ───── SLIDE 17: ROTOR CONFIGS ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "COVERAGE")
    add_page_num(s, pn)
    add_title(s, "F135 Rotor Coverage: 9 Stages")
    add_subtitle(s, "12 part numbers, 20-110 blades per stage")

    stage_data = [
        ["1", "20", "4137321, 4135411, 4134621"],
        ["2", "38", "4130812"],
        ["3 (focus)", "55", "4134613, 4130813"],
        ["4", "42", "4119904"],
        ["5-6", "54-70", "4119905, 4133006"],
        ["7-9", "82-110", "4136007, 4133008, 4131129-01"],
    ]
    st = s.shapes.add_table(7, 3, Inches(0.5), Inches(1.4), Inches(9), Inches(2.6))
    stable = st.table
    for ci, h in enumerate(["Stage", "Blades", "Part Numbers"]):
        cell = stable.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(11)
                r.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(stage_data):
        for ci, val in enumerate(row):
            cell = stable.cell(ri + 1, ci)
            cell.text = val
            if ri % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TINT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(11)
                    r.font.name = "Arial"
                    r.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    add_insight_bar(s, "SCALABILITY", "Zone boundaries stored as percentages, not absolute mm. One config works across all 9 stages with different blade geometries.")

    # ───── SLIDE 18: SECTION DIVIDER - DASHBOARD ─────
    s = dark_slide(pres)
    pn = sn()
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "04"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    tf2 = txBox2.text_frame
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "Interactive Dashboard"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Georgia"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    run3 = tf3.paragraphs[0].add_run()
    run3.text = "FastAPI web application with Plotly 3D visualization and live reporting"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    run3.font.name = "Arial"

    # ───── SLIDE 19: DASHBOARD FEATURES ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "DASHBOARD")
    add_page_num(s, pn)
    add_title(s, "Web Dashboard: Three Views")
    add_subtitle(s, "FastAPI + Plotly for interactive inspection and reporting")

    add_card(s, "Dashboard (Main)", "Stats cards, 3D deviation heatmap, disposition pie chart,\nzone distribution, defect table, depth comparison bar chart.\nAll interactive. All real-time data from pipeline output.", 0.5, 1.4, w=4.2, h=1.5)
    add_card(s, "3D Viewer", "Full-screen rotatable point cloud. Four view modes:\nDeviation Heatmap, Scan, CAD Reference, Defect Clusters.\nPoint size slider, max points selector.", 5.2, 1.4, w=4.2, h=1.5)
    add_card(s, "Reports Page", "All inspection reports with timestamps and dispositions.\nJSON viewer inline. Excel and JSON download buttons.\nOne-click pipeline re-run from the dashboard.", 0.5, 3.1, w=4.2, h=0.9)
    add_card(s, "REST API", "12 endpoints: /api/report/latest, /api/pointcloud/deviations,\n/api/pipeline/run, /api/config, and more.\nReady for integration with other Hitachi systems.", 5.2, 3.1, w=4.2, h=0.9)

    add_insight_bar(s, "LIVE DEMO", "The dashboard is running right now at localhost:8000. We will switch to a live demo after this slide.")

    # ───── SLIDE 20: AWS ARCHITECTURE ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "CLOUD")
    add_page_num(s, pn)
    add_title(s, "Production Architecture: AWS")
    add_subtitle(s, "Parallel GPU processing for <90s per IBR at any blade count")

    add_body_text(s, "S3 Upload \u2192 EventBridge \u2192 Lambda Orchestrator \u2192 AWS Batch (parallel)\n\nAll N blades processed simultaneously on separate g4dn.xlarge GPUs.\nEach blade: ~20 seconds. Total wall time: 20s + 40s aggregation = 60s.\n\nStage 3 (55 blades): $0.18 per IBR\nStage 9 (110 blades): $0.34 per IBR\n\nITAR compliant: all processing in AWS US regions, encrypted at rest,\nVPC with no internet gateway, CloudTrail audit logging.", x=0.5, y=1.4, h=2.8, size=12)

    add_insight_bar(s, "PRODUCTION PATH", "Architecture is designed. Sprint 4 will provision the AWS infrastructure and deploy the containerized pipeline.")

    # ───── SLIDE 21: SECTION DIVIDER - NEXT STEPS ─────
    s = dark_slide(pres)
    pn = sn()
    txBox = s.shapes.add_textbox(Inches(0.8), Inches(0.3), Inches(1), Inches(1))
    tf = txBox.text_frame
    run = tf.paragraphs[0].add_run()
    run.text = "05"
    run.font.size = Pt(60)
    run.font.bold = True
    run.font.color.rgb = PRIMARY_LIGHT
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.4), Inches(0.6))
    tf2 = txBox2.text_frame
    run2 = tf2.paragraphs[0].add_run()
    run2.text = "Next Steps & Sprint 4"
    run2.font.size = Pt(32)
    run2.font.bold = True
    run2.font.color.rgb = WHITE
    run2.font.name = "Georgia"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.3), Inches(4), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    txBox3 = s.shapes.add_textbox(Inches(0.8), Inches(2.6), Inches(8.4), Inches(0.4))
    tf3 = txBox3.text_frame
    run3 = tf3.paragraphs[0].add_run()
    run3.text = "What needs to happen before production readiness"
    run3.font.size = Pt(14)
    run3.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    run3.font.name = "Arial"

    # ───── SLIDE 22: OPEN QUESTIONS ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "QUESTIONS FOR P&W")
    add_page_num(s, pn)
    add_title(s, "Open Questions for Pratt & Whitney")
    add_subtitle(s, "Items we need resolved to move forward")

    questions = [
        ("Zone boundaries by stage", "Do zones A1-A3, B1-B2 scale with blade size per stage, or are they fixed?"),
        ("LE/TE curve data", "We need actual leading/trailing edge curves from the CAD for edge-distance classification."),
        ("Real scan data", "Can we get a sample Stage 3 PLY scan for validation against vendor results?"),
        ("Adjacency rules", "Spanwise-only distance or Euclidean distance for cross-defect adjacency checks?"),
        ("Cross-foil limits", "Are there assembly-level limits beyond individual blade rules?"),
    ]
    for i, (title, desc) in enumerate(questions):
        add_card(s, title, desc, 0.5, 1.4 + i * 0.52, w=9, h=0.46)

    add_insight_bar(s, "BLOCKING ITEMS", "Real scan data is the single biggest blocker. Everything else we can work around with reasonable assumptions.")

    # ───── SLIDE 23: SPRINT 4 PLAN ─────
    s = light_slide(pres)
    pn = sn()
    add_badge(s, "SPRINT 4")
    add_page_num(s, pn)
    add_title(s, "Sprint 4 Priorities")
    add_subtitle(s, "Two-week plan to move from prototype to validation-ready")

    sprint_items = [
        ("AWS Infrastructure", "Provision RDS, Lambda, Batch compute environment, S3 buckets."),
        ("Real Scan Validation", "Run pipeline against actual P&W scan data. Compare to vendor results."),
        ("LE/TE Curve Extraction", "Extract edge curves from CAD mesh for accurate edge-distance measurement."),
        ("ML Model Training", "Train Random Forest on labeled defect data for defect type classification."),
        ("Foil Segmentation Tuning", "Calibrate DBSCAN parameters on real rotor geometry (not synthetic)."),
    ]
    for i, (title, desc) in enumerate(sprint_items):
        add_card(s, f"{i+1}. {title}", desc, 0.5, 1.4 + i * 0.52, w=9, h=0.46)

    add_insight_bar(s, "SPRINT 4 GOAL", "By end of Sprint 4, run the pipeline against a real Stage 3 scan and produce a report that matches vendor findings within tolerance.")

    # ───── SLIDE 24: CLOSING (DARK) ─────
    s = dark_slide(pres)
    pn = sn()

    txBox = s.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(0.8))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Pipeline is Running. Results are Real."
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = WHITE
    run.font.name = "Georgia"

    txBox2 = s.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(8.4), Inches(0.4))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = "14.5 seconds. 4 defects detected. 91 rules checked. Reports generated."
    run2.font.size = Pt(14)
    run2.font.color.rgb = PRIMARY_LIGHT
    run2.font.name = "Arial"

    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.6), Inches(3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = PRIMARY_LIGHT
    line.line.fill.background()

    for i, txt in enumerate([
        "Supreeth Mysore  |  supreeth.mysore@hitachids.com",
        "Hitachi Digital Services  |  Pratt & Whitney F135 Program",
        "github.com/drmysore/open3d-ibr-defect-extraction",
    ]):
        txBox = s.shapes.add_textbox(Inches(0.8), Inches(3.0 + i * 0.35), Inches(8.4), Inches(0.3))
        tf = txBox.text_frame
        run = tf.paragraphs[0].add_run()
        run.text = txt
        run.font.size = Pt(12)
        run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
        run.font.name = "Arial"

    output_path = os.path.join("output", "IBR_Defect_Extraction_Sprint3_Review.pptx")
    os.makedirs("output", exist_ok=True)
    pres.save(output_path)
    print(f"Presentation saved: {output_path}")
    return output_path


if __name__ == "__main__":
    build_presentation()
